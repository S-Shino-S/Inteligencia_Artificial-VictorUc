#!/usr/bin/env python3
"""Generate a short Wumpus-world lecture deck (AIMA fig. 7.2)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

W = Inches(13.333333)
H = Inches(7.5)
NAVY_DARK = RGBColor(0x0A, 0x1F, 0x40)
NAVY = RGBColor(0x0F, 0x2C, 0x59)
ACCENT = RGBColor(0x1B, 0x6C, 0xA8)
GOLD = RGBColor(0xE8, 0xB9, 0x3A)
SLATE = RGBColor(0x2C, 0x3E, 0x50)
MUTED = RGBColor(0x5D, 0x6D, 0x7E)
PAPER = RGBColor(0xF6, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD5, 0xDE, 0xE8)
SUB = RGBColor(0xC5, 0xD0, 0xDC)

FOOTER = "Curso de Inteligencia Artificial  ·  Agentes"
BASE = Path("/Users/victoruccetina/Documents/code/inteligencia-artificial/Agentes")
OUT = BASE / "PPTXs" / "wumpus-world.pptx"
IMG = BASE / "images" / "Wumpus world.png"
TOTAL = 11


def set_run(run, text, size, color, bold=False, name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def set_anchor(tf, anchor):
    mapping = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
    tf._txBody.bodyPr.set("anchor", mapping[anchor])


def add_rect(slide, l, t, w, h, fill, line=None, line_w=Pt(1)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = line_w
    return sh


def add_round(slide, l, t, w, h, fill, line=LINE, line_w=Pt(1.25)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = line_w
    try:
        sh.adjustments[0] = 0.1
    except Exception:
        pass
    return sh


def add_tb(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size, color, bold)
    return box


def add_paras(slide, l, t, w, h, items, size=16, color=SLATE, space_after=10, bullet=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    set_anchor(tf, MSO_ANCHOR.TOP)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        run = p.add_run()
        set_run(run, ("•  " if bullet else "") + item, size, color, False)
    return box


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def chrome(slide, n):
    add_rect(slide, 0, 0, W, H, WHITE)
    add_rect(slide, 0, 0, Inches(0.12), H, GOLD)
    add_rect(slide, 0, Inches(7.12), W, Inches(0.38), NAVY)
    add_tb(slide, Inches(0.5), Inches(7.14), Inches(10.0), Inches(0.32), FOOTER, 11, WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_tb(
        slide, Inches(11.35), Inches(7.14), Inches(1.5), Inches(0.32),
        f"{n} / {TOTAL}", 11, WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )


def heading(slide, title, subtitle=None):
    add_tb(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.55), title, 28, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.55), Inches(0.88), Inches(1.4), Inches(0.06), GOLD)
    if subtitle:
        add_tb(slide, Inches(0.55), Inches(1.02), Inches(12.2), Inches(0.4), subtitle, 16, MUTED, anchor=MSO_ANCHOR.MIDDLE)


def card(slide, l, t, w, h, kicker, title, body, body_size=14):
    add_round(slide, l, t, w, h, PAPER)
    add_tb(slide, l + Inches(0.22), t + Inches(0.18), w - Inches(0.44), Inches(0.32), kicker, 13, ACCENT, bold=True)
    add_tb(slide, l + Inches(0.22), t + Inches(0.48), w - Inches(0.44), Inches(0.55), title, 18, NAVY, bold=True)
    add_tb(slide, l + Inches(0.22), t + Inches(1.08), w - Inches(0.44), h - Inches(1.28), body, body_size, SLATE)


def pic(slide, path, l, t, w=None, h=None):
    kwargs = {}
    if w is not None:
        kwargs["width"] = w
    if h is not None:
        kwargs["height"] = h
    slide.shapes.add_picture(str(path), l, t, **kwargs)


def title_slide(prs):
    s = blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_DARK)
    add_rect(s, 0, 0, Inches(0.18), H, GOLD)
    add_tb(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.4), "CURSO DE INTELIGENCIA ARTIFICIAL", 16, GOLD, bold=True)
    add_tb(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.2), "El mundo del Wumpus", 40, WHITE, bold=True)
    add_tb(
        s, Inches(0.9), Inches(4.0), Inches(11.0), Inches(1.2),
        "Una cueva 4×4 que el agente no ve entera:\npozos, un monstruo, oro y cinco perceptos.",
        18, SUB,
    )
    add_tb(
        s, Inches(0.9), Inches(6.55), Inches(11.0), Inches(0.35),
        "AIMA, capítulos 2 y 7  ·  figura 7.2",
        14, GOLD,
    )


def agenda(prs):
    s = blank(prs)
    chrome(s, 2)
    heading(s, "Agenda")
    items = [
        ("01", "La cueva", "Cuadrícula 4×4. (1, 1) abajo a la izquierda. El mapa de la figura es la vista de Dios."),
        ("02", "Peligros y oro", "Wumpus, pozos y una barra de oro. Matar o caer termina el episodio."),
        ("03", "Perceptos y acciones", "Cinco bits de sensor. Seis acciones. No hay un canal extra de «intención»."),
        ("04", "La tarea", "Agarrar el oro, volver a la salida y trepar. Cada paso cuesta."),
    ]
    positions = [
        (Inches(0.55), Inches(1.55)), (Inches(6.95), Inches(1.55)),
        (Inches(0.55), Inches(4.2)), (Inches(6.95), Inches(4.2)),
    ]
    for item, pos in zip(items, positions):
        card(s, pos[0], pos[1], Inches(5.85), Inches(2.35), *item)


def figure_slide(prs):
    s = blank(prs)
    chrome(s, 3)
    heading(s, "La figura 7.2", "Esto es el mundo real. El agente, al empezar, solo está en (1, 1) y no ve el resto")
    pic(s, IMG, Inches(3.95), Inches(1.45), h=Inches(5.2))


def coords_slide(prs):
    s = blank(prs)
    chrome(s, 4)
    heading(s, "Cómo se lee el tablero", "Columnas a la derecha, filas hacia arriba. Igual que AIMA")
    pic(s, IMG, Inches(0.4), Inches(1.45), h=Inches(5.2))
    add_round(s, Inches(6.85), Inches(1.5), Inches(5.85), Inches(5.1), PAPER)
    add_paras(
        s, Inches(7.1), Inches(1.75), Inches(5.4), Inches(4.6),
        [
            "Una celda es una cueva (x, y), 1-based.",
            "(1, 1) es la esquina inferior izquierda: ahí nace el agente.",
            "Empieza mirando al este (hacia la columna 2).",
            "Adyacente = norte, sur, este u oeste. No hay diagonales.",
            "Chocar con el muro te deja donde estás y da Bump.",
        ],
        size=16, space_after=10, bullet=True,
    )


def hazards_slide(prs):
    s = blank(prs)
    chrome(s, 5)
    heading(s, "Tres peligros y un premio", "En esta figura: Wumpus en (1, 3), pozos en (3, 1), (3, 3) y (4, 4), oro en (2, 3)")
    cards = [
        ("Wumpus", "Un monstruo. Si entras a su cueva y sigue vivo, te come. En las cuevas de al lado hay hedor (Stench)."),
        ("Pozo", "Caer es morir. En las cuevas de al lado hay brisa (Breeze). Varios pozos a la vez: las brisas se superponen."),
        ("Oro", "Una barra. En esa cueva brilla (Glitter). Hay que agarrarla: no se pega sola a los pies."),
        ("Salida", "Solo se trepa (Climb) desde (1, 1). Si no volviste ahí, no sales."),
    ]
    positions = [
        (Inches(0.55), Inches(1.55)), (Inches(6.95), Inches(1.55)),
        (Inches(0.55), Inches(4.2)), (Inches(6.95), Inches(4.2)),
    ]
    for (t, b), (l, top) in zip(cards, positions):
        add_round(s, l, top, Inches(5.85), Inches(2.35), PAPER)
        add_tb(s, l + Inches(0.28), top + Inches(0.28), Inches(5.3), Inches(0.5), t, 18, NAVY, bold=True)
        add_tb(s, l + Inches(0.28), top + Inches(0.9), Inches(5.3), Inches(1.2), b, 15, SLATE)


def percepts_slide(prs):
    s = blank(prs)
    chrome(s, 6)
    heading(s, "El percepto es una 5-tupla", "Cada paso el entorno devuelve [Stench, Breeze, Glitter, Bump, Scream]. Nada más")
    items = [
        ("Stench", "Hedor: un vecino (ortogonal) tiene al Wumpus vivo."),
        ("Breeze", "Brisa: un vecino es un pozo."),
        ("Glitter", "Brillo: el oro está en esta cueva (aún no lo has agarrado)."),
        ("Bump", "Tope: el Forward chocó con el borde. No te moviste."),
        ("Scream", "Grito: tu flecha mató al Wumpus (en la línea de tiro)."),
        ("[None]", "Los cinco bits en falso: (1, 1) en la figura. Ni brisa ni hedor."),
    ]
    for i, (t, b) in enumerate(items):
        col, row = i % 3, i // 3
        l = Inches(0.55) + Inches(col * 4.2)
        top = Inches(1.55) + Inches(row * 2.7)
        add_round(s, l, top, Inches(3.95), Inches(2.5), PAPER)
        add_tb(s, l + Inches(0.22), top + Inches(0.22), Inches(3.5), Inches(0.5), t, 18, NAVY, bold=True)
        add_tb(s, l + Inches(0.22), top + Inches(0.8), Inches(3.5), Inches(1.45), b, 15, SLATE)


def actions_slide(prs):
    s = blank(prs)
    chrome(s, 7)
    heading(s, "Seis acciones", "El agente elige una. El mundo responde con el percepto nuevo y un cambio de puntuación")
    items = [
        ("Forward", "Un paso en la dirección a la que miras. Si hay muro: Bump."),
        ("TurnLeft / TurnRight", "Giras 90°. No cambias de cueva. Cuesta un paso igual."),
        ("Grab", "Si hay Glitter aquí, recoges el oro. Si no, no pasa nada."),
        ("Shoot", "Gasta la única flecha en línea recta. Si da al Wumpus: Scream."),
        ("Climb", "Solo en (1, 1). Termina el episodio. Con oro, ganas; sin oro, sales vacío."),
        ("Una por turno", "No hay «andar y agarrar» a la vez. Grab y luego Forward, o al revés."),
    ]
    for i, (t, b) in enumerate(items):
        col, row = i % 3, i // 3
        l = Inches(0.55) + Inches(col * 4.2)
        top = Inches(1.55) + Inches(row * 2.7)
        add_round(s, l, top, Inches(3.95), Inches(2.5), PAPER)
        add_tb(s, l + Inches(0.22), top + Inches(0.22), Inches(3.5), Inches(0.55), t, 17, NAVY, bold=True)
        add_tb(s, l + Inches(0.22), top + Inches(0.85), Inches(3.5), Inches(1.4), b, 15, SLATE)


def score_slide(prs):
    s = blank(prs)
    chrome(s, 8)
    heading(s, "La puntuación", "El episodio suma. Un camino corto con oro gana a un paseo largo con oro")
    cards = [
        ("+1000", "Sales por (1, 1) con el oro (Climb). No basta con estar encima de la barra."),
        ("−1000", "Pozo o Wumpus. El episodio termina. Da igual el oro que llevaras."),
        ("−1", "Cada acción (incluido girar). Empuja a no dar vueltas."),
        ("−10", "Disparar. Una flecha: piénsalo. Fallar no mata al Wumpus."),
    ]
    positions = [
        (Inches(0.55), Inches(1.55)), (Inches(6.95), Inches(1.55)),
        (Inches(0.55), Inches(4.2)), (Inches(6.95), Inches(4.2)),
    ]
    for (t, b), (l, top) in zip(cards, positions):
        add_round(s, l, top, Inches(5.85), Inches(2.35), PAPER)
        add_tb(s, l + Inches(0.28), top + Inches(0.28), Inches(5.3), Inches(0.5), t, 22, NAVY, bold=True)
        add_tb(s, l + Inches(0.28), top + Inches(0.95), Inches(5.3), Inches(1.15), b, 16, SLATE)


def read_figure(prs):
    s = blank(prs)
    chrome(s, 9)
    heading(s, "Cómo se lee la figura", "El agente no tiene este dibujo. Nosotros sí: sirve para comprobar una deducción")
    pic(s, IMG, Inches(0.4), Inches(1.45), h=Inches(5.2))
    add_round(s, Inches(6.85), Inches(1.5), Inches(5.85), Inches(5.1), PAPER)
    add_paras(
        s, Inches(7.1), Inches(1.75), Inches(5.4), Inches(4.6),
        [
            "(1, 1): ni brisa ni hedor. (1, 2) y (2, 1) son seguras.",
            "(2, 1): brisa. Hay un pozo en un vecino: aquí es (3, 1).",
            "(1, 2): hedor. El Wumpus está al lado: aquí, (1, 3).",
            "Tras ver (2, 1) y (1, 2), (2, 2) no puede ser pozo ni Wumpus: es el atajo seguro.",
            "El oro está en (2, 3), con brisa y hedor: vecinos peligrosos, la celda misma no es pozo.",
        ],
        size=15, space_after=9, bullet=True,
    )


def task_slide(prs):
    s = blank(prs)
    chrome(s, 10)
    heading(s, "La tarea", "No es «encontrar el oro». Es salir vivo con el oro")
    cards = [
        ("1. Ir", "Avanzar solo por cuevas que ya puedes justificar como seguras."),
        ("2. Grab", "En (2, 3) hay Glitter. Sin Grab, el oro se queda. Climb vacío no paga +1000."),
        ("3. Volver", "El Climb solo vale en (1, 1). Hay que deshacer el camino (o uno equivalente)."),
        ("4. Climb", "Termina. El marcador es 1000 menos pasos y flechas. Morir vale −1000 más los pasos."),
    ]
    positions = [
        (Inches(0.55), Inches(1.55)), (Inches(6.95), Inches(1.55)),
        (Inches(0.55), Inches(4.2)), (Inches(6.95), Inches(4.2)),
    ]
    for (t, b), (l, top) in zip(cards, positions):
        add_round(s, l, top, Inches(5.85), Inches(2.35), PAPER)
        add_tb(s, l + Inches(0.28), top + Inches(0.28), Inches(5.3), Inches(0.5), t, 18, NAVY, bold=True)
        add_tb(s, l + Inches(0.28), top + Inches(0.9), Inches(5.3), Inches(1.2), b, 16, SLATE)


def takeaways(prs):
    s = blank(prs)
    chrome(s, 11)
    heading(s, "Ideas para llevar", "El Wumpus es un mundo parcialmente observable con un sensor muy pobre")
    items = [
        ("Vista de Dios", "La figura 7.2 es para nosotros. El programa solo recibe la 5-tupla de ahora."),
        ("Local", "Stench y Breeze hablan de vecinos, no de la cueva en la que estás."),
        ("Memoria", "Sin recordar (2, 1) y (1, 2) no puedes declarar (2, 2) segura. Un percepto no basta."),
        ("Objetivo", "Oro + salida. Un agente que camina al azar suele caer al pozo de (3, 1)."),
    ]
    positions = [
        (Inches(0.55), Inches(1.55)), (Inches(6.95), Inches(1.55)),
        (Inches(0.55), Inches(4.2)), (Inches(6.95), Inches(4.2)),
    ]
    for (t, b), (l, top) in zip(items, positions):
        add_round(s, l, top, Inches(5.85), Inches(2.35), PAPER)
        add_tb(s, l + Inches(0.28), top + Inches(0.28), Inches(5.3), Inches(0.5), t, 18, NAVY, bold=True)
        add_tb(s, l + Inches(0.28), top + Inches(0.9), Inches(5.3), Inches(1.2), b, 15, SLATE)


def main():
    if not IMG.exists():
        raise SystemExit(f"missing image: {IMG}")
    (BASE / "PPTXs").mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    title_slide(prs)
    agenda(prs)
    figure_slide(prs)
    coords_slide(prs)
    hazards_slide(prs)
    percepts_slide(prs)
    actions_slide(prs)
    score_slide(prs)
    read_figure(prs)
    task_slide(prs)
    takeaways(prs)
    prs.save(OUT)
    print(f"Wrote {len(prs.slides)} slides → {OUT}")


if __name__ == "__main__":
    main()
